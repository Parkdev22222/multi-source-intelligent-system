#!/usr/bin/env python3
"""Render the paper's figure from docs/method_figure.pptx.

The PowerPoint file is the source of the diagram: edit a box there and the
paper picks it up, rather than the two drifting. This reads the shape tree back
out and re-emits it as the vector PDF the paper includes.

    python scripts/figure_from_pptx.py     # docs/method_figure.{pdf,png}

Two things are deliberately not carried over. The slide's title is dropped --
in the paper the caption does that job, and a heading inside a float reads as a
second one. And the slide's margins go with it: the output is cropped to the
drawing, so \\columnwidth is spent on the diagram instead of on whitespace.

Calibri is not installed here, so the text is set in Liberation Sans, which is
metric-compatible with Arial rather than with Calibri: the .pptx and this PDF
will not be glyph-identical. The geometry, colours and strings are.
"""
from __future__ import annotations

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
matplotlib.rcParams["pdf.fonttype"] = 42        # TrueType, never Type 3
matplotlib.rcParams["ps.fonttype"] = 42
import matplotlib.pyplot as plt
from matplotlib.patches import Ellipse, FancyArrowPatch, FancyBboxPatch
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "docs" / "method_figure.pptx"
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
EMU = 914400.0
FONT = "Liberation Sans"
DROP = "Vision-grounded change reports"      # the slide title, not the figure's

# The paper's column is 3.5in; the drawing keeps the aspect it has on the slide.
FIG_W = 3.5


def rgb(obj, default=None):
    try:
        return "#" + str(obj.rgb)
    except Exception:
        return default


def main() -> int:
    if not SRC.is_file():
        print(f"missing {SRC.relative_to(ROOT)}", file=sys.stderr)
        return 2
    shapes = list(Presentation(SRC).slides[0].shapes)
    shapes = [s for s in shapes if not (s.has_text_frame and DROP in s.text_frame.text)]

    # Crop to the ink: the slide's margins are not the figure's.
    xs = [s.left / EMU for s in shapes] + [(s.left + s.width) / EMU for s in shapes]
    ys = [s.top / EMU for s in shapes] + [(s.top + s.height) / EMU for s in shapes]
    x0, x1, y0, y1 = min(xs), max(xs), min(ys), max(ys)
    pad = 0.04
    x0, x1, y0, y1 = x0 - pad, x1 + pad, y0 - pad, y1 + pad
    scale = FIG_W / (x1 - x0)

    fig, ax = plt.subplots(figsize=(FIG_W, (y1 - y0) * scale))
    ax.set_xlim(x0, x1); ax.set_ylim(y1, y0); ax.axis("off")
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)

    for sh in shapes:
        x, y = sh.left / EMU, sh.top / EMU
        w, h = sh.width / EMU, sh.height / EMU
        geom = sh.element.find(f".//{A}prstGeom")
        prst = geom.get("prst") if geom is not None else "rect"
        ln = sh.element.spPr.find(f"{A}ln") if sh.element.spPr is not None else None
        edge = rgb(sh.line.color) if ln is not None else None
        lw = (sh.line.width / 12700.0) * scale if sh.line.width else 0.8
        try:
            fill = rgb(sh.fill.fore_color) if sh.fill.type == 1 else None
        except Exception:
            fill = None

        if prst == "line":
            xf = sh.element.spPr.xfrm
            fh, fv = xf.get("flipH") == "1", xf.get("flipV") == "1"
            xa, xb = (x + w, x) if fh else (x, x + w)
            ya, yb = (y + h, y) if fv else (y, y + h)
            dash = ln.find(f"{A}prstDash") if ln is not None else None
            style = {"sysDot": (0, (1, 1.6)), "dash": (0, (2.5, 2))}.get(
                dash.get("val") if dash is not None else "", "-")
            head = ln is not None and ln.find(f"{A}tailEnd") is not None
            if head:
                ax.add_patch(FancyArrowPatch((xa, ya), (xb, yb), arrowstyle="-|>",
                                             mutation_scale=7, linewidth=lw,
                                             color=edge or "#888", linestyle=style,
                                             shrinkA=0, shrinkB=0, zorder=4))
            else:
                ax.plot([xa, xb], [ya, yb], color=edge or "#888", linewidth=lw,
                        linestyle=style, solid_capstyle="round", zorder=3)
            continue

        if prst == "ellipse":
            ax.add_patch(Ellipse((x + w/2, y + h/2), w, h, facecolor=fill or "#888",
                                 edgecolor=edge or "none", zorder=6))
        elif fill or edge:
            # A box can be dashed too -- Change-Fact-Score is, to mark it as the
            # evaluator rather than a pipeline stage. Reading prstDash only for
            # lines silently squared it off into an ordinary box.
            bdash = ln.find(f"{A}prstDash") if ln is not None else None
            bstyle = {"sysDot": (0, (1, 1.6)), "dash": (0, (2.5, 2))}.get(
                bdash.get("val") if bdash is not None else "", "-")
            ax.add_patch(FancyBboxPatch((x, y), w, h,
                                        boxstyle="round,pad=0,rounding_size=0.05",
                                        facecolor=fill or "none",
                                        edgecolor=edge or "none", linewidth=lw,
                                        linestyle=bstyle, zorder=2))

        if sh.has_text_frame and sh.text_frame.text.strip():
            rot = 90 if sh.rotation else 0
            # One paragraph is one line, and each carries its own size: a box
            # holds a label over one or two smaller sub-labels. Stacking has to
            # walk paragraphs, not runs -- indexing a run list by line number
            # put every line at the same height.
            lines = []
            for para in sh.text_frame.paragraphs:
                if not para.text.strip():
                    continue
                r = para.runs[0]
                lines.append((para.text,
                              (r.font.size.pt if r.font.size else 12) * scale,
                              rgb(r.font.color, "#1a1d23") or "#1a1d23",
                              bool(r.font.italic), bool(r.font.bold)))
            if not lines:
                continue
            # The axes are in slide inches, not figure inches: fs has already
            # been scaled down for the 3.5in column, so dividing it back out
            # gives the line's height in the space the boxes are drawn in.
            heights = [fs / 72.0 / scale * 1.32 for _, fs, _, _, _ in lines]
            cy = y + h / 2 - sum(heights) / 2
            for (text, fs, color, italic, bold), lh in zip(lines, heights):
                cy += lh / 2
                ax.text(x + w / 2, cy, text, ha="center", va="center",
                        fontsize=fs, color=color, rotation=rot, fontfamily=FONT,
                        zorder=7, style="italic" if italic else "normal",
                        weight="bold" if bold else "normal")
                cy += lh / 2

    for out in (ROOT / "docs" / "method_figure.pdf", ROOT / "docs" / "method_figure.png"):
        fig.savefig(out, dpi=400, bbox_inches="tight", pad_inches=0.01)
        print(f"wrote {out.relative_to(ROOT)}")
    plt.close(fig)
    return 0


if __name__ == "__main__":
    sys.exit(main())
